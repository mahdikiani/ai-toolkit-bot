import asyncio
import logging
from io import BytesIO

from server.config import Settings
from utils import dify, finance, imagetools, mime, pdftools, texttools

from .models import OcrTask

CONVERTING_IMAGE_EXTS = {"image/png", "image/tiff"}
IMAGE_EXTS = {"image/jpeg"}
DOCX_EXTS = {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
PDF_EXTS = {"application/pdf"}
PPTX_EXTS = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
}


async def process_ocr_task_service(task: OcrTask) -> OcrTask:
    logging.info("Starting processing for task %s", task.id)
    file_content = await task.file_content()
    file_type = mime.check_file_type(file_content)
    if file_type in DOCX_EXTS:
        res = extract_docx(file_content)
        return await save_result(task, res)
    elif file_type in PPTX_EXTS:
        res = extract_pptx(file_content)
        return await save_result(task, res)
    elif file_type in PDF_EXTS:
        pages = pdftools.extract_pdf_bytes_pages(file_content)
    elif file_type in CONVERTING_IMAGE_EXTS:
        pages = [imagetools.convert_to_jpg_bytes(file_content)]
    elif file_type in IMAGE_EXTS:
        pages = [file_content]
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    quota = await finance.check_quota(task.user_id, len(pages), raise_exception=False)
    if quota < len(pages):
        task.task_status = "error"
        await task.save_report("insufficient_quota")
        return task

    async with dify.AsyncDifyClient(Settings.pishrun_api_key) as client:
        sem = asyncio.Semaphore(16)

        async def sem_ocr(page: BytesIO) -> str | None:
            async with sem:
                return await client.ocr_image(page)

        text_pages = await asyncio.gather(*(sem_ocr(page) for page in pages))

    usage = await finance.meter_cost(task.user_id, len(pages))
    await save_result(
        task,
        "\n\n".join([t for t in text_pages if t]),
        usage_amount=usage.amount,
        usage_id=usage.uid,
    )

    return task


async def save_result(
    task: OcrTask,
    result: str,
    usage_amount: float | None = None,
    usage_id: str | None = None,
) -> OcrTask:
    task.result = texttools.normalize_text(result)
    task.task_status = "completed"
    task.usage_amount = usage_amount
    task.usage_id = usage_id
    return await task.save()


def extract_docx(file_content: BytesIO) -> str:
    import docx

    doc = docx.Document(file_content)
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text)
    return texttools.normalize_text(text)


def extract_pptx(file_content: BytesIO) -> str:
    import pptx

    prs = pptx.Presentation(file_content)
    texts = []
    for slide in prs.slides:
        texts.extend(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    return texttools.normalize_text("\n\n".join(texts))
